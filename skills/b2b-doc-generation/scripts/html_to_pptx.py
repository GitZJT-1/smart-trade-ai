#!/usr/bin/env python3
"""
html_to_pptx.py — Convert an existing HTML brochure to editable PPTX.
Usage:  uv run python html_to_pptx.py <input.html> [output.pptx]

Maps HTML page divs (cover / inner / back-cover) to PowerPoint slides
with base64 image decoding, styled shapes, and native text boxes.

Dependencies: python-pptx, beautifulsoup4, lxml
"""

import base64, io, os, re, sys
from bs4 import BeautifulSoup

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ─── Config ──────────────────────────────────────────────────────
if len(sys.argv) < 2:
    print("Usage: uv run python html_to_pptx.py <input.html> [output.pptx]")
    sys.exit(1)

HTML_PATH = sys.argv[1]
if len(sys.argv) >= 3:
    OUTPUT_PATH = sys.argv[2]
else:
    OUTPUT_PATH = os.path.splitext(HTML_PATH)[0] + ".pptx"


# ─── Helpers ──────────────────────────────────────────────────────
def decode_b64(src):
    if not src or not src.startswith("data:"):
        return None, None
    m = re.match(r'data:([^;]+);base64,(.+)', src)
    if not m:
        return None, None
    return m.group(1), base64.b64decode(m.group(2))

def add_image(slide, b64_src, left, top, width, height=None):
    try:
        mime, data = decode_b64(b64_src)
        if data is None:
            return None
        return slide.shapes.add_picture(io.BytesIO(data), left, top, width, height or width)
    except Exception:
        return None

def add_shape_bg(slide, color_hex, left=0, top=0,
                  width=Inches(13.333), height=Inches(7.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*parse_hex(color_hex))
    shape.line.fill.background()
    return shape

def parse_hex(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color="#333333",
                 align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*parse_hex(color))
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_rich_text(slide, lines, left, top, width, height,
                  default_size=11, default_color="#444444"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, sz, bd, cl, al = item, default_size, False, default_color, PP_ALIGN.LEFT
        else:
            txt = item[0]
            sz = item[1] if len(item) > 1 else default_size
            bd = item[2] if len(item) > 2 else False
            cl = item[3] if len(item) > 3 else default_color
            al = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(txt)
        p.font.size = Pt(sz)
        p.font.bold = bd
        p.font.color.rgb = RGBColor(*parse_hex(cl))
        p.alignment = al
    return txBox

def add_gold_bar(slide, left, top, width=Inches(1.5), height=Pt(3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(232, 196, 74)
    shape.line.fill.background()
    return shape


# ─── Read HTML ────────────────────────────────────────────────────
with open(HTML_PATH, "r", encoding="utf-8") as f:
    html = f.read()

soup = BeautifulSoup(html, "lxml")

# ─── Create Presentation ──────────────────────────────────────────
prs = Presentation()
SLIDE_W = prs.slide_width = Inches(13.333)
SLIDE_H = prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# ─── Extract pages — direct body children only! ───────────────────
body = soup.find("body")
pages = []
if body:
    for div in body.find_all("div", recursive=False):
        classes = div.get("class", [])
        if any(c in ("page", "cover", "inner", "back-cover") for c in classes):
            pages.append(div)

print(f"Found {len(pages)} page elements")

# ─── Build slides ─────────────────────────────────────────────────
for idx, page_div in enumerate(pages):
    classes = page_div.get("class", [])
    cls_str = " ".join(classes)
    slide = prs.slides.add_slide(blank_layout)
    imgs = page_div.find_all("img")
    all_text = page_div.get_text(separator="\n", strip=True)

    is_cover = "cover" in cls_str and "back-cover" not in cls_str
    is_backcover = "back-cover" in cls_str
    is_inner = "inner" in cls_str

    # ── Cover ──────────────────────────────────────────────────
    if is_cover:
        add_shape_bg(slide, "#0f2027")
        for img in imgs:
            src = img.get("src", "")
            if src.startswith("data:image"):
                add_image(slide, src, Inches(9.5), 0, Inches(3.8), SLIDE_H)
                break
        h1 = page_div.find("h1")
        h2 = page_div.find("h2")
        tagline = page_div.find(class_="tagline")
        tagline_en = page_div.find(class_="tagline-en")
        add_text_box(slide, h1.get_text(strip=True) if h1 else "",
                     Inches(1), Inches(1.8), Inches(8), Inches(0.8),
                     font_size=40, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)
        add_text_box(slide, h2.get_text(strip=True) if h2 else "",
                     Inches(1), Inches(2.6), Inches(8), Inches(0.5),
                     font_size=16, color="#CCCCCC", align=PP_ALIGN.CENTER)
        add_gold_bar(slide, Inches(4.5), Inches(3.2))
        if tagline:
            add_text_box(slide, tagline.get_text(strip=True),
                         Inches(1), Inches(3.5), Inches(8), Inches(0.5),
                         font_size=22, bold=True, color="#e8c44a", align=PP_ALIGN.CENTER)
        if tagline_en:
            add_text_box(slide, tagline_en.get_text(strip=True),
                         Inches(1), Inches(4.0), Inches(8), Inches(0.4),
                         font_size=12, color="#999999", align=PP_ALIGN.CENTER)
        bottom_div = page_div.find(class_="bottom")
        if bottom_div:
            add_text_box(slide, bottom_div.get_text(strip=True),
                         Inches(1), Inches(6.5), Inches(8), Inches(0.3),
                         font_size=10, color="#777777", align=PP_ALIGN.CENTER)

    # ── Back Cover ────────────────────────────────────────────
    elif is_backcover:
        add_shape_bg(slide, "#2c5364")
        deco_div = page_div.find(class_="deco-section")
        if deco_div:
            d_img = deco_div.find("img")
            if d_img and d_img.get("src", "").startswith("data:image"):
                add_image(slide, d_img["src"], Inches(9.5), 0, Inches(3.8), SLIDE_H)
        add_text_box(slide, "沈阳百欧通用机械有限公司",
                     Inches(1), Inches(1.5), Inches(8), Inches(0.6),
                     font_size=28, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)
        add_text_box(slide, "SHENYANG BAI OU GENERAL MACHINERY CO., LTD.",
                     Inches(1), Inches(2.1), Inches(8), Inches(0.4),
                     font_size=12, color="#B0B0B0", align=PP_ALIGN.CENTER)
        add_gold_bar(slide, Inches(4.5), Inches(2.7))
        info_div = page_div.find(class_="info")
        if info_div:
            lines = info_div.get_text("\n", strip=True).split("\n")
            add_rich_text(slide, [(l, 11, False, "#DDDDDD", PP_ALIGN.CENTER) for l in lines],
                          Inches(1), Inches(3.0), Inches(8), Inches(2.5))
        copyright_div = page_div.find(class_="copyright")
        if copyright_div:
            add_text_box(slide, copyright_div.get_text(strip=True),
                         Inches(1), Inches(6.5), Inches(8), Inches(0.5),
                         font_size=9, color="#888888", align=PP_ALIGN.CENTER)

    # ── Inner Pages ───────────────────────────────────────────
    elif is_inner:
        add_shape_bg(slide, "#FFFFFF")
        header = page_div.find(class_="page-header")
        if header:
            st = header.find(class_="section-title")
            ste = header.find(class_="section-title-en")
            if st:
                add_text_box(slide, st.get_text(strip=True),
                             Inches(0.5), Inches(0.3), Inches(8), Inches(0.5),
                             font_size=24, bold=True, color="#1a3a4a")
            if ste:
                add_text_box(slide, ste.get_text(strip=True),
                             Inches(0.5), Inches(0.7), Inches(8), Inches(0.3),
                             font_size=10, color="#8a9ba8")
            add_gold_bar(slide, Inches(0.5), Inches(0.9), Inches(0.6), Pt(3))

        # Data bar
        y_off = 1.1
        data_bar = page_div.find(class_="data-bar")
        if data_bar:
            items = data_bar.find_all(class_="data-item")
            for i, item in enumerate(items):
                num = item.find(class_="num")
                label = item.find(class_="label")
                ntext = num.get_text(strip=True) if num else ""
                ltext = label.get_text(strip=True) if label else ""
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.5 + i * 2.5), Inches(1.1), Inches(2.3), Inches(0.7))
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(26, 58, 74)
                shape.line.fill.background()
                add_text_box(slide, ntext, Inches(0.5 + i * 2.5), Inches(1.12),
                             Inches(2.3), Inches(0.3),
                             font_size=18, bold=True, color="#e8c44a", align=PP_ALIGN.CENTER)
                add_text_box(slide, ltext, Inches(0.5 + i * 2.5), Inches(1.4),
                             Inches(2.3), Inches(0.3),
                             font_size=8, color="#CCCCCC", align=PP_ALIGN.CENTER)
            y_off = 1.9

        # Image grid
        img_grid = page_div.find(class_="img-grid")
        if img_grid:
            cards = img_grid.find_all(class_="img-card")
            for k, card in enumerate(cards):
                ci, ri = k % 3, k // 3
                card_img = card.find("img")
                if card_img and card_img.get("src", "").startswith("data:image"):
                    add_image(slide, card_img["src"],
                              Inches(0.5 + ci * 4.2), Inches(y_off + ri * 1.6),
                              Inches(3.8), Inches(1.2))
                label = card.find(class_="label")
                if label:
                    add_text_box(slide, label.get_text(strip=True),
                        Inches(0.5 + ci * 4.2), Inches(y_off + ri * 1.6 + 1.2),
                        Inches(3.8), Inches(0.3),
                        font_size=9, bold=True, color="#333333", align=PP_ALIGN.CENTER)

        # Page number
        if idx <= 10:
            add_text_box(slide, f"{idx:02d}/10",
                         Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3),
                         font_size=8, color="#CCCCCC", align=PP_ALIGN.RIGHT)


# ─── Save ─────────────────────────────────────────────────────────
prs.save(OUTPUT_PATH)
fsize = os.path.getsize(OUTPUT_PATH)
print(f"✅ PPTX saved: {OUTPUT_PATH}")
print(f"   Size: {fsize/1024:.1f} KB  |  Slides: {len(prs.slides)}")
